from pptx import Presentation
from pptx.util import Emu

prs = Presentation("/Users/gorkemcetin/Desktop/Casbah.pptx")
print(f"Slide size: {prs.slide_width} x {prs.slide_height}")
print(f"  In inches: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f}")
print(f"Slide count: {len(prs.slides)}")
print()

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} (layout: {slide.slide_layout.name}) ---")
    for shape in slide.shapes:
        kind = shape.shape_type
        name = shape.name
        if shape.has_text_frame:
            text = " | ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())[:200]
            fonts = set()
            sizes = set()
            colors = set()
            bolds = set()
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name: fonts.add(r.font.name)
                    if r.font.size: sizes.add(r.font.size.pt)
                    if r.font.bold is not None: bolds.add(r.font.bold)
                    try:
                        if r.font.color and r.font.color.rgb:
                            colors.add(str(r.font.color.rgb))
                    except Exception:
                        pass
            print(f"  [{name}] text={text!r}")
            if fonts: print(f"    fonts={fonts} sizes={sizes} bolds={bolds} colors={colors}")
        else:
            print(f"  [{name}] kind={kind}")
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    rgb = fill.fore_color.rgb
                    print(f"    fill={rgb}")
                except Exception:
                    pass
        except Exception:
            pass
    print()
