from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

GREEN = RGBColor(0x13, 0x71, 0x5B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD5, 0xDD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


s = prs.slides.add_slide(BLANK)

add_rect(s, 0, 0, SW, Inches(0.08), GREEN)
add_text(s, Inches(0.6), Inches(0.25), Inches(10), Inches(0.4),
         "WHAT GARTNER IS SAYING", size=12, bold=True, color=GREEN)
add_text(s, Inches(0.6), Inches(0.55), Inches(12), Inches(0.7),
         "Why VerifyWise. Why now.", size=32, bold=True, color=DARK)
add_rect(s, Inches(0.6), Inches(1.3), Inches(1.0), Inches(0.05), GREEN)

stats = [
    ("3.4x", "more effective at AI governance",
        "Organisations that deploy an AI governance platform are more than three times as effective at governance than those that don't."),
    ("13%", "of leaders feel ready to govern AI",
        "Only 13% of IT and business leaders feel fully equipped to lead AI governance today."),
    ("40%", "fewer AI ethical incidents by 2028",
        "Enterprises that adopt a comprehensive AI governance platform will see 40% fewer AI-related ethical incidents."),
    ("50%", "of enterprises on a converged platform by 2030",
        "More than half of enterprises will run on a single platform that converges data, analytics and AI governance."),
]

x0, y0 = Inches(0.6), Inches(1.55)
tw, th = Inches(6.0), Inches(2.45)
gap_x, gap_y = Inches(0.15), Inches(0.2)
for i, (big, lead, body) in enumerate(stats):
    col = i % 2
    row = i // 2
    x = x0 + (tw + gap_x) * col
    y = y0 + (th + gap_y) * row
    add_rect(s, x, y, tw, th, LIGHT, line=BORDER)
    add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(2.4), Inches(1.1),
             big, size=52, bold=True, color=GREEN)
    add_text(s, x + Inches(2.8), y + Inches(0.45), tw - Inches(3.0), Inches(0.7),
             lead, size=15, bold=True, color=DARK)
    add_text(s, x + Inches(0.3), y + Inches(1.45), tw - Inches(0.5), Inches(0.95),
             body, size=13, color=GREY)

add_rect(s, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.55), GREEN)
add_text(s, Inches(0.8), Inches(6.72), Inches(11.7), Inches(0.4),
         "Gartner: invest in an AI governance platform before AI risks become impossible to manage with existing processes.",
         size=14, bold=True, color=WHITE)

add_text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
         "Sources: Gartner 2025 State of AI-Ready Data Survey; 2026 Chief Data and Analytics Officer Agenda Survey; Generative and Agentic AI in Enterprise Applications Survey.",
         size=9, color=GREY)

out = "/Users/gorkemcetin/Desktop/VerifyWise_Why_Now_Gartner.pptx"
prs.save(out)
print(out)
