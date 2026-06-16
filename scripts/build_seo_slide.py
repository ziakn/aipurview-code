from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

GREEN = RGBColor(0x13, 0x71, 0x5B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD5, 0xDD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "• " + it
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


s = prs.slides.add_slide(BLANK)

# Header bar + title
add_rect(s, 0, 0, SW, Inches(0.08), GREEN)
add_text(s, Inches(0.6), Inches(0.25), Inches(10), Inches(0.4),
         "INDEPENDENT VISIBILITY AUDIT, JUNE 2026", size=12, bold=True, color=GREEN)
add_text(s, Inches(0.6), Inches(0.55), Inches(12), Inches(0.7),
         "Defining the category on Google and inside AI assistants.",
         size=28, bold=True, color=DARK)
add_rect(s, Inches(0.6), Inches(1.3), Inches(1.0), Inches(0.05), GREEN)

# Top row: three stat tiles
tiles = [
    ("#1", "on Google for\n\"Open Source AI\nGovernance Tool\""),
    ("#3", "on Google for\n\"EU AI Act\nGovernance Platform\""),
    ("4 / 4", "AI models name\nVerifyWise on open-source\nand community queries"),
]
x = Inches(0.6)
for big, body in tiles:
    add_rect(s, x, Inches(1.55), Inches(4.0), Inches(2.3), LIGHT, line=BORDER)
    add_text(s, x + Inches(0.3), Inches(1.7), Inches(3.6), Inches(1.0),
             big, size=52, bold=True, color=GREEN)
    add_text(s, x + Inches(0.3), Inches(2.7), Inches(3.6), Inches(1.1),
             body, size=14, color=DARK)
    x += Inches(4.2)

# Bottom row: two panels (winning / emerging)
add_rect(s, Inches(0.6), Inches(4.0), Inches(6.1), Inches(2.7), WHITE, line=BORDER)
add_text(s, Inches(0.8), Inches(4.1), Inches(5.7), Inches(0.4),
         "WHERE WE ARE WINNING", size=12, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(4.45), Inches(5.7), Inches(2.2), [
    "#1 on Google for open-source AI governance",
    "#1 for ISO 42001 + EU AI Act long-tail searches",
    "Recognised by all 4 AI models on open-source queries",
    "Strong feature and framework-specific SEO coverage",
], size=14)

add_rect(s, Inches(6.85), Inches(4.0), Inches(5.85), Inches(2.7), WHITE, line=BORDER)
add_text(s, Inches(7.05), Inches(4.1), Inches(5.5), Inches(0.4),
         "WHERE WE ARE EMERGING", size=12, bold=True, color=GREEN)
add_bullets(s, Inches(7.05), Inches(4.45), Inches(5.5), Inches(2.2), [
    "Top 5 on Google for LLM governance and self-hosted queries",
    "Gemini and ChatGPT cite VerifyWise on NIST AI RMF + ISO 42001",
    "Perplexity and Gemini cite us on free-tier and scaling queries",
    "Buyer-intent recognition now in 2 of the 4 models",
], size=14)

# Footer thesis bar
add_rect(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5), GREEN)
add_text(s, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
         "VerifyWise isn't fighting for visibility. It's defining the open-source, self-hosted, EU AI Act-ready category.",
         size=13, bold=True, color=WHITE)

out = "/Users/gorkemcetin/Desktop/VerifyWise_SEO_AI_Visibility.pptx"
prs.save(out)
print(out)
