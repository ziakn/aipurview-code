from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

GREEN = RGBColor(0x13, 0x71, 0x5B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD5, 0xDD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = "• " + lead
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = GREEN
            r2 = p.add_run()
            r2.text = ". " + rest
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_header(slide, title, kicker=None):
    add_rect(slide, 0, 0, SW, Inches(0.08), GREEN)
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.25), Inches(8), Inches(0.4),
                 kicker.upper(), size=12, bold=True, color=GREEN)
        add_text(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.7),
                 title, size=30, bold=True, color=DARK)
    else:
        add_text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8),
                 title, size=30, bold=True, color=DARK)
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(1.0), Inches(0.05), GREEN)


def add_footer(slide, page):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "VerifyWise and NTT DATA  |  Confidential", size=10, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{page} / 10", size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ---------- SLIDE 1: COVER ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, Inches(4.6), SH, GREEN)
add_text(s, Inches(0.6), Inches(0.6), Inches(3.5), Inches(0.5),
         "VERIFYWISE", size=14, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.0), Inches(3.8), Inches(1.2),
         "AI governance,\nmade operational.", size=34, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(4.6), Inches(3.8), Inches(0.5),
         "A partnership opportunity for NTT DATA", size=16, color=WHITE)

add_text(s, Inches(5.2), Inches(0.8), Inches(7.5), Inches(0.4),
         "PARTNERSHIP BRIEFING", size=12, bold=True, color=GREEN)
add_text(s, Inches(5.2), Inches(1.2), Inches(7.5), Inches(2.0),
         "From advisory hours\nto recurring software\nrevenue.",
         size=36, bold=True, color=DARK)
add_rect(s, Inches(5.2), Inches(3.6), Inches(0.8), Inches(0.05), GREEN)
add_text(s, Inches(5.2), Inches(3.75), Inches(7.5), Inches(0.4),
         "Prepared for NTT DATA sales leadership", size=16, bold=True, color=DARK)
add_text(s, Inches(5.2), Inches(4.15), Inches(7.5), Inches(0.4),
         "Ulas Ozguven, Co-founder, VerifyWise", size=14, color=GREY)
add_text(s, Inches(5.2), Inches(4.5), Inches(7.5), Inches(0.4),
         "June 2026", size=14, color=GREY)

add_rect(s, Inches(5.2), Inches(6.0), Inches(7.5), Inches(0.9), LIGHT)
add_text(s, Inches(5.4), Inches(6.15), Inches(7.3), Inches(0.6),
         "Assess. Govern. Prove. As a managed service.",
         size=16, bold=True, color=GREEN)

# ---------- SLIDE 2: THE OPPORTUNITY ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "The 2026 deadline window", kicker="The opportunity")
add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
         "Clients are deploying AI faster than they can prove it's safe.",
         size=20, bold=True, color=DARK)

tiles = [
    ("€35M", "Maximum EU AI Act fine,\nor a percentage of global\nturnover. Phased through\n2026."),
    ("100s", "Smart AI Agents NTT DATA\nhas already deployed\ninto regulated sectors."),
    ("0", "Branded governance\nplatforms inside NTT\nDATA's stack today."),
]
x = Inches(0.6)
for headline, body in tiles:
    add_rect(s, x, Inches(2.3), Inches(4.0), Inches(2.6), LIGHT, line=BORDER)
    add_text(s, x + Inches(0.3), Inches(2.5), Inches(3.6), Inches(0.9),
             headline, size=44, bold=True, color=GREEN)
    add_text(s, x + Inches(0.3), Inches(3.5), Inches(3.6), Inches(1.4),
             body, size=15, color=DARK)
    x += Inches(4.2)

add_rect(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.5), WHITE, line=GREEN)
add_text(s, Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.4),
         "THE GAP", size=12, bold=True, color=GREEN)
add_text(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.0),
         "NTT DATA sells AI transformation. Clients now need proof: inventory, risk, "
         "evidence, audit. Today that proof lives in spreadsheets. VerifyWise turns it into a product.",
         size=16, color=DARK)
add_footer(s, 2)

# ---------- SLIDE 3: WHERE NTT DATA IS TODAY ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "NTT DATA's AI governance position today", kicker="Current state")

left_items = [
    ("Strong advisory practice", "GRC consulting and trustworthy-AI guidance for boards."),
    ("Hundreds of AI agents in production", "Healthcare, automotive, finance. All high-risk under the EU AI Act."),
    ("Sovereign cloud launch partner", "AWS European Sovereign Cloud, for strict data residency clients."),
]
right_items = [
    ("No software layer", "Advice points clients at generic tools. No branded platform underneath."),
    ("Liability grows with every agent", "Each deployment adds a control surface the client has to defend."),
    ("Shadow AI blind spot", "Current monitoring covers threats and behaviour. Not unsanctioned AI."),
]

add_text(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.5),
         "Strengths", size=20, bold=True, color=GREEN)
add_bullets(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), left_items, size=15)

add_text(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(0.5),
         "Gaps", size=20, bold=True, color=DARK)
add_bullets(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(4.8), right_items, size=15)

add_rect(s, Inches(6.65), Inches(1.6), Inches(0.02), Inches(4.8), BORDER)
add_footer(s, 3)

# ---------- SLIDE 4: WHAT VERIFYWISE BRINGS ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "What VerifyWise brings", kicker="The platform")

add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
         "The software layer your advisory practice doesn't have yet.",
         size=18, color=DARK)

capabilities = [
    ("AI system inventory", "Every model, vendor and use case in one register."),
    ("Risk register and FRIA", "Lines up with EU AI Act Article 9 risk management."),
    ("Multi-framework engine", "EU AI Act, ISO 42001, ISO 27001, NIST AI RMF, plus 30 plugins."),
    ("Evidence hub", "Upload once, satisfy many controls. Audit-ready."),
    ("AI Gateway and guardrails", "PII detection, spend caps and policy checks at runtime."),
    ("Shadow AI detection", "Find unsanctioned AI use across the enterprise."),
]

x0, y0 = Inches(0.6), Inches(2.3)
tw, th = Inches(4.0), Inches(1.4)
gap_x, gap_y = Inches(0.15), Inches(0.2)
for i, (title, body) in enumerate(capabilities):
    col = i % 3
    row = i // 3
    x = x0 + (tw + gap_x) * col
    y = y0 + (th + gap_y) * row
    add_rect(s, x, y, tw, th, WHITE, line=BORDER)
    add_text(s, x + Inches(0.25), y + Inches(0.15), tw - Inches(0.35), Inches(0.4),
             title, size=15, bold=True, color=DARK)
    add_text(s, x + Inches(0.25), y + Inches(0.55), tw - Inches(0.35), Inches(0.8),
             body, size=13, color=GREY)

add_text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.6),
         "Self-hosted or SaaS. Multi-tenant. Built for regulated industries.",
         size=15, bold=True, color=GREEN)
add_footer(s, 4)

# ---------- SLIDE 5: STRATEGIC FIT (5 PILLARS) ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Five places this partnership pays off", kicker="Strategic fit")

pillars = [
    ("1", "Productise the advisory practice",
        "Anchor GRC consulting in a branded platform. Higher margin and easier to repeat across accounts."),
    ("2", "Hit the EU AI Act deadline",
        "European clients have obligations live in 2026. They need something they can deploy, not a slide deck."),
    ("3", "Wrap the agents NTT DATA already shipped",
        "Hundreds of Smart AI Agents in high-risk sectors. Each one is a governance surface waiting to be proven."),
    ("4", "Sovereign cloud fit",
        "VerifyWise self-hosts cleanly inside NTT DATA's managed sovereign cloud environments."),
    ("5", "Shadow AI as a wedge",
        "A detection capability the other SIs don't currently offer. Useful in almost every first conversation."),
]

y = Inches(1.5)
for num, title, body in pillars:
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.9), WHITE, line=BORDER)
    add_rect(s, Inches(0.6), y, Inches(0.9), Inches(0.9), GREEN)
    add_text(s, Inches(0.6), y + Inches(0.18), Inches(0.9), Inches(0.6),
             num, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), y + Inches(0.12), Inches(10.8), Inches(0.4),
             title, size=17, bold=True, color=DARK)
    add_text(s, Inches(1.7), y + Inches(0.45), Inches(10.8), Inches(0.5),
             body, size=14, color=GREY)
    y += Inches(1.05)

add_footer(s, 5)

# ---------- SLIDE 6: MANAGED SERVICE PLAY ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "AI governance as a managed service", kicker="The commercial play")

add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
         "One-off advisory becomes a subscription with services around it.",
         size=18, bold=True, color=DARK)

stages = [
    ("DISCOVER", "Inventory AI\nsystems and\nshadow AI"),
    ("ASSESS", "Risk class,\nFRIA, gap\nanalysis"),
    ("GOVERN", "Controls,\npolicies,\nguardrails"),
    ("PROVE", "Evidence,\naudit,\nreporting"),
    ("MONITOR", "Ongoing\noversight and\nincidents"),
]
x = Inches(0.6)
w = Inches(2.3)
gap = Inches(0.15)
for i, (label, body) in enumerate(stages):
    add_rect(s, x, Inches(2.3), w, Inches(2.0), GREEN if i % 2 == 0 else DARK)
    add_text(s, x + Inches(0.15), Inches(2.45), w - Inches(0.3), Inches(0.5),
             label, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.0), w - Inches(0.3), Inches(1.3),
             body, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w, Inches(3.05), gap, Inches(0.5))
        arr.fill.solid()
        arr.fill.fore_color.rgb = GREEN
        arr.line.fill.background()
    x += w + gap

add_rect(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.0), LIGHT, line=BORDER)
add_text(s, Inches(0.85), Inches(4.85), Inches(11.5), Inches(0.5),
         "REVENUE MODEL", size=13, bold=True, color=GREEN)

rev_items = [
    ("Platform subscription", "Recurring license per client. Annuity revenue rather than project-based."),
    ("Managed service fees", "NTT DATA runs governance operations on the client's behalf. Healthy margin."),
    ("Follow-on consulting", "Policy work, ISO 42001 implementation, audits, training, model evaluations."),
]
add_bullets(s, Inches(0.85), Inches(5.2), Inches(11.5), Inches(1.5), rev_items, size=14)
add_footer(s, 6)

# ---------- SLIDE 7: CLIENT OUTCOMES ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "What NTT DATA's clients get out of it", kicker="Client outcomes")

outcomes = [
    ("Faster AI to production", "Intake, risk and approval workflows unblock the legal and security gates."),
    ("Audit-ready by default", "Evidence is collected as work happens, rather than assembled in a panic."),
    ("Board-level visibility", "One dashboard for compliance posture across every framework in scope."),
    ("Lower regulatory risk", "EU AI Act, ISO 42001 and NIST AI RMF positions you can actually prove."),
    ("Vendor and model assurance", "You know which third-party model is powering which use case."),
    ("Catch shadow AI early", "Spot unsanctioned tools before they show up in an incident report."),
]

y = Inches(1.5)
for i, (title, body) in enumerate(outcomes):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.2)
    yy = y + row * Inches(1.65)
    add_rect(s, x, yy, Inches(6.0), Inches(1.45), WHITE, line=BORDER)
    add_text(s, x + Inches(0.3), yy + Inches(0.2), Inches(5.6), Inches(0.5),
             title, size=17, bold=True, color=DARK)
    add_text(s, x + Inches(0.3), yy + Inches(0.7), Inches(5.6), Inches(0.7),
             body, size=14, color=GREY)
add_footer(s, 7)

# ---------- SLIDE 8: DIFFERENTIATION ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Where VerifyWise wins", kicker="Differentiation")

cols = ["Capability", "Spreadsheets / GRC", "Generic AI tool", "VerifyWise"]
rows = [
    ("AI system inventory",       "Manual",          "Partial",      "Native"),
    ("EU AI Act and ISO 42001",   "DIY",             "One framework", "5 plus 30 plugins"),
    ("Shadow AI detection",       "None",            "None",         "Built in"),
    ("AI Gateway guardrails",     "None",            "Partial",      "Runtime PII and policy"),
    ("Self-host / sovereign",     "None",            "Rare",         "Yes"),
    ("Audit evidence hub",        "Folders",         "Light",        "Reusable across frameworks"),
]

table_x, table_y = Inches(0.6), Inches(1.5)
col_widths = [Inches(3.4), Inches(2.9), Inches(2.9), Inches(2.9)]
row_h = Inches(0.55)

x = table_x
for i, c in enumerate(cols):
    fill = GREEN if i == 3 else DARK
    add_rect(s, x, table_y, col_widths[i], row_h, fill)
    add_text(s, x + Inches(0.15), table_y + Inches(0.12), col_widths[i] - Inches(0.2), Inches(0.4),
             c, size=14, bold=True, color=WHITE)
    x += col_widths[i]

for r, row in enumerate(rows):
    yy = table_y + row_h + r * row_h
    x = table_x
    bg = LIGHT if r % 2 == 0 else WHITE
    for i, val in enumerate(row):
        add_rect(s, x, yy, col_widths[i], row_h, bg, line=BORDER)
        is_vw = i == 3
        add_text(s, x + Inches(0.15), yy + Inches(0.12), col_widths[i] - Inches(0.2), Inches(0.4),
                 val, size=13, bold=is_vw, color=GREEN if is_vw else DARK)
        x += col_widths[i]

add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
         "Plenty of SIs can build AI. Fewer can help a client prove it. That's the gap we close.",
         size=16, bold=True, color=GREEN)
add_footer(s, 8)

# ---------- SLIDE 9: PARTNERSHIP MODEL ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "How we partner", kicker="Engagement model")

tiers = [
    ("REFERRAL", "Lightweight",
        ["Co-sell into existing accounts",
         "Standard referral economics",
         "Joint case studies"]),
    ("RESELL AND DELIVER", "Recommended",
        ["NTT DATA-branded managed service",
         "Tiered platform margins",
         "Delivery enablement and training",
         "Joint go-to-market across EU and UK"]),
    ("STRATEGIC OEM", "Long-term",
        ["White-label option for sovereign cloud",
         "Joint input on the product roadmap",
         "Co-investment in regulated verticals"]),
]
x = Inches(0.6)
w = Inches(4.0)
for i, (name, tag, items) in enumerate(tiers):
    highlight = i == 1
    add_rect(s, x, Inches(1.5), w, Inches(5.0), GREEN if highlight else WHITE, line=GREEN if highlight else BORDER)
    add_text(s, x + Inches(0.25), Inches(1.7), w - Inches(0.5), Inches(0.5),
             name, size=18, bold=True, color=WHITE if highlight else DARK)
    add_text(s, x + Inches(0.25), Inches(2.15), w - Inches(0.5), Inches(0.4),
             tag, size=12, bold=True, color=WHITE if highlight else GREEN)
    add_rect(s, x + Inches(0.25), Inches(2.6), Inches(0.6), Inches(0.04),
             WHITE if highlight else GREEN)
    tb = s.shapes.add_textbox(x + Inches(0.25), Inches(2.8), w - Inches(0.5), Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = "• " + it
        r.font.name = "Calibri"
        r.font.size = Pt(15)
        r.font.color.rgb = WHITE if highlight else DARK
    x += w + Inches(0.15)

add_footer(s, 9)

# ---------- SLIDE 10: NEXT STEPS ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Let's pick the first three accounts", kicker="Next steps")

steps = [
    ("Week 1-2", "Identify three NTT DATA accounts with live EU AI Act exposure. VerifyWise joins the discovery calls."),
    ("Week 3-6", "Run paid pilots. Co-branded, time-boxed, with success criteria agreed up front."),
    ("Week 7-10", "Convert pilots to subscription plus managed service. Build the first joint case study."),
    ("Quarter 2", "Formalise the partnership tier. Joint EU and UK go-to-market. Expand into the next verticals."),
]
y = Inches(1.5)
for label, body in steps:
    add_rect(s, Inches(0.6), y, Inches(2.2), Inches(1.1), GREEN)
    add_text(s, Inches(0.7), y + Inches(0.3), Inches(2.0), Inches(0.5),
             label, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(2.85), y, Inches(9.85), Inches(1.1), WHITE, line=BORDER)
    add_text(s, Inches(3.05), y + Inches(0.2), Inches(9.5), Inches(0.85),
             body, size=15, color=DARK)
    y += Inches(1.2)

add_rect(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.7), GREEN)
add_text(s, Inches(0.85), Inches(6.42), Inches(11.7), Inches(0.5),
         "Ulas Ozguven  |  Co-founder, VerifyWise  |  ulas@verifywise.ai",
         size=16, bold=True, color=WHITE)
add_footer(s, 10)

out = "/Users/gorkemcetin/Desktop/VerifyWise_NTT_DATA_Partnership.pptx"
prs.save(out)
print(out)
