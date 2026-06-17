"""
VerifyWise capabilities deck for Asociación Española Contra el Cáncer (AECC).
Contact: Juan Pedro Benítez. Written by: Ulas Ozdemir, VerifyWise.
English version first; Spanish to follow.
Style: VerifyWise design guide (green #13715B, off-white, card grids, kicker headers).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

GREEN = RGBColor(0x13, 0x71, 0x5B)
AMBER = RGBColor(0xE8, 0xA0, 0x20)  # brand v1.2 accent — stats/keywords only, 2-3x max
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD5, 0xDD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

TOTAL = 9


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = "• " + lead
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = GREEN
            r2 = p.add_run()
            r2.text = " — " + rest
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_header(slide, title, kicker=None):
    # Minimal: thin green top bar + title only. No kicker label, no accent bar.
    add_rect(slide, 0, 0, SW, Inches(0.08), GREEN)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.8),
             title, size=30, bold=True, color=DARK)


LOGO = "/Users/gorkemcetin/verifywise/verifywise-logo-full.png"
LOGO_W = Inches(1.5)
LOGO_H = Inches(1.5) * (226 / 1200)  # keep aspect ratio


def add_footer(slide, page):
    # Small VerifyWise logo bottom-left, page number bottom-right. Nothing else.
    slide.shapes.add_picture(LOGO, Inches(0.6), Inches(7.02), width=LOGO_W, height=LOGO_H)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.8), Inches(0.3),
             str(page), size=10, color=GREY, align=PP_ALIGN.RIGHT)


def card_grid(slide, items, x0, y0, cols, tw, th, gap_x, gap_y, title_size=15, body_size=13):
    for i, (title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = x0 + (tw + gap_x) * col
        y = y0 + (th + gap_y) * row
        add_rect(slide, x, y, tw, th, WHITE, line=BORDER)
        add_text(slide, x + Inches(0.22), y + Inches(0.15), tw - Inches(0.35), Inches(0.45),
                 title, size=title_size, bold=True, color=DARK)
        add_text(slide, x + Inches(0.22), y + Inches(0.58), tw - Inches(0.35), th - Inches(0.6),
                 body, size=body_size, color=GREY)


# ---------- SLIDE 1: COVER ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, Inches(4.6), SH, GREEN)
add_text(s, Inches(0.6), Inches(0.6), Inches(3.5), Inches(0.5),
         "VERIFYWISE", size=14, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.0), Inches(3.8), Inches(1.6),
         "AI governance,\nrisk and\ncompliance.", size=32, bold=True, color=WHITE)

add_text(s, Inches(5.2), Inches(1.9), Inches(7.5), Inches(0.4),
         "Prepared for", size=15, color=GREY)
add_text(s, Inches(5.2), Inches(2.35), Inches(7.5), Inches(1.4),
         "Asociación Española\nContra el Cáncer", size=34, bold=True, color=DARK)
add_text(s, Inches(5.2), Inches(4.05), Inches(7.5), Inches(0.4),
         "Juan Pedro Benítez", size=18, bold=True, color=DARK)
add_text(s, Inches(5.2), Inches(5.0), Inches(7.5), Inches(0.4),
         "Prepared by", size=15, color=GREY)
add_text(s, Inches(5.2), Inches(5.4), Inches(7.5), Inches(0.4),
         "Ulas Ozdemir, VerifyWise", size=18, bold=True, color=DARK)

# ---------- SLIDE 2: WHO WE ARE ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "AI risk, governance and compliance", kicker="Who we are")
add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.0),
         "VerifyWise is an AI governance platform. It helps organisations track their AI "
         "systems, score the risk, put controls in place and keep the evidence an auditor will "
         "ask for. We are a UK company, built for European regulation.",
         size=18, color=DARK)

tiles = [
    ("UK company", "Headquartered in the UK,\nbuilt for the EU AI Act and\nEuropean compliance."),
    ("16+ modules", "Covering the full lifecycle\nof AI governance, from intake\nto public disclosure."),
    ("24+ frameworks", "EU AI Act, ISO 42001,\nISO 27001, NIST AI RMF,\nGDPR, SOC 2 and more."),
]
x = Inches(0.6)
for headline, body in tiles:
    add_rect(s, x, Inches(2.9), Inches(4.0), Inches(2.5), LIGHT, line=BORDER)
    add_text(s, x + Inches(0.3), Inches(3.1), Inches(3.6), Inches(0.7),
             headline, size=26, bold=True, color=AMBER)
    add_text(s, x + Inches(0.3), Inches(3.9), Inches(3.6), Inches(1.4),
             body, size=15, color=DARK)
    x += Inches(4.2)

add_rect(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.0), WHITE, line=GREEN)
add_text(s, Inches(0.85), Inches(5.85), Inches(11.7), Inches(0.7),
         "Governance you can actually deploy and run, not a policy document that sits in a drawer.",
         size=16, bold=True, color=GREEN)
add_footer(s, 2)

# ---------- SLIDE 3: PRODUCT CAPABILITIES (LLM Gateway) ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Product capabilities", kicker="What the platform does")
add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5),
         "One control point in front of your AI traffic. Observability, compliance and cost, on every request.",
         size=17, color=DARK)

caps = [
    ("Observability", "See every model request: who called what, latency, spend and any guardrail that fired."),
    ("Compliance", "Gateway activity maps to EU AI Act articles, so it doubles as audit evidence."),
    ("Cost management", "Per-team budgets and spend tracking. Requests stop when the budget runs out."),
    ("Guardrails", "In-process PII detection and content filters. Block or mask before it reaches a model."),
    ("Virtual keys", "A revocable key per team, with its own budget and rate limit. No shared credentials."),
    ("Agent governance", "The same controls, approvals and audit trail for autonomous agents over MCP."),
]
card_grid(s, caps, Inches(0.6), Inches(2.2), 3, Inches(4.0), Inches(1.5), Inches(0.15), Inches(0.2))

add_text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.6),
         "Provider keys encrypted at rest. 100+ providers behind one OpenAI-compatible endpoint.",
         size=15, bold=True, color=GREEN)
add_footer(s, 3)

# ---------- SLIDE 4: DEPLOYMENT OPTIONS ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Deployment options", kicker="How you run it")
add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5),
         "Run it our way or yours. Self-hosted deployments get the same features, updates and support.",
         size=17, color=DARK)

deploys = [
    ("SaaS", "Our cloud. Teams start in\nminutes, no infrastructure\nto manage."),
    ("On-premise", "Your infrastructure, via\nDocker or Kubernetes.\nRunning in days."),
    ("Custom on-premise", "Tailored to your\nenvironment, data residency\nand security requirements."),
]
x = Inches(0.6)
for headline, body in deploys:
    add_rect(s, x, Inches(2.3), Inches(4.0), Inches(2.4), WHITE, line=BORDER)
    add_rect(s, x, Inches(2.3), Inches(4.0), Inches(0.7), GREEN)
    add_text(s, x + Inches(0.3), Inches(2.42), Inches(3.5), Inches(0.5),
             headline, size=20, bold=True, color=WHITE)
    add_text(s, x + Inches(0.3), Inches(3.2), Inches(3.5), Inches(1.4),
             body, size=15, color=DARK)
    x += Inches(4.2)

add_rect(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.6), LIGHT, line=BORDER)
add_text(s, Inches(0.85), Inches(5.25), Inches(11.7), Inches(0.4),
         "SOURCE-AVAILABLE, NO LOCK-IN", size=12, bold=True, color=GREEN)
add_bullets(s, Inches(0.85), Inches(5.6), Inches(11.7), Inches(1.1), [
    ("Inspectable code", "your security team can audit the platform before you deploy it."),
    ("Data sovereignty", "your data stays where you decide. Not an add-on."),
    ("Export everything", "your audit trails and documentation are yours. If you leave, you take it all."),
], size=14)
add_footer(s, 4)

# ---------- SLIDE 5: DELIVERY MODELS ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Delivery models", kicker="How we engage")

# Software only
add_rect(s, Inches(0.6), Inches(1.5), Inches(5.9), Inches(5.0), WHITE, line=BORDER)
add_rect(s, Inches(0.6), Inches(1.5), Inches(5.9), Inches(0.8), DARK)
add_text(s, Inches(0.85), Inches(1.66), Inches(5.4), Inches(0.5),
         "Software only", size=20, bold=True, color=WHITE)
add_bullets(s, Inches(0.85), Inches(2.6), Inches(5.4), Inches(3.6), [
    ("SaaS", "managed for you on our cloud."),
    ("On-premise", "run it inside your own environment."),
    ("Custom on-premise", "tailored to your security and data residency needs."),
], size=16)

# Software + consulting
add_rect(s, Inches(6.85), Inches(1.5), Inches(5.85), Inches(5.0), GREEN)
add_text(s, Inches(7.1), Inches(1.66), Inches(5.4), Inches(0.5),
         "Software + consulting", size=20, bold=True, color=WHITE)
tb = s.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.4), Inches(3.8))
tf = tb.text_frame
tf.word_wrap = True
plus = [
    ("Fractional AI governance officer", "an experienced governance lead embedded with your team, part-time."),
    ("Deploy, run, transition", "we stand the platform up, operate it, then hand the keys back when you are ready."),
    ("Customisation as a service", "we adapt VerifyWise to how your teams already work."),
]
for j, (lead, rest) in enumerate(plus):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    r1 = p.add_run(); r1.text = "• " + lead
    r1.font.name = "Calibri"; r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = " — " + rest
    r2.font.name = "Calibri"; r2.font.size = Pt(15); r2.font.color.rgb = WHITE
add_footer(s, 5)

# ---------- SLIDE 6: PROFESSIONAL SERVICES ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "VerifyWise professional services", kicker="Beyond the software")
add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5),
         "We do not just hand over a login. We help you set governance up and keep it running.",
         size=17, color=DARK)

svc = [
    ("Deploy", "We stand up VerifyWise in your environment, configured for your frameworks and teams."),
    ("Run", "We operate governance on your behalf: monitoring, evidence collection and reporting."),
    ("Transfer", "We train your team and hand operations over to you, with no lock-in either way."),
    ("Customise", "We adapt the platform to how your teams already work."),
    ("Spanish language support", "The platform is available in Spanish, so your teams work in their own language."),
    ("European focus", "Built around the EU AI Act and the regulations European organisations answer to."),
]
card_grid(s, svc, Inches(0.6), Inches(2.2), 3, Inches(4.0), Inches(1.6), Inches(0.15), Inches(0.2))
add_footer(s, 6)

# ---------- SLIDE 7: SPANISH SUPPORT + CLOSING ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, SW, Inches(0.08), GREEN)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(1.0),
         "Available in Spanish,\nready for Europe.", size=36, bold=True, color=DARK)

add_bullets(s, Inches(0.6), Inches(3.0), Inches(12), Inches(2.2), [
    ("Full Spanish interface", "your team uses VerifyWise entirely in Spanish."),
    ("EU AI Act native", "risk classification, record-keeping and oversight mapped to the regulation."),
    ("Self-hosted or SaaS", "with the same features and support either way."),
    ("Software or software plus services", "take the platform alone, or with our team alongside it."),
], size=17)

add_rect(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.9), GREEN)
add_text(s, Inches(0.85), Inches(6.02), Inches(11.7), Inches(0.6),
         "Ulas Ozdemir  |  VerifyWise  |  verifywise.ai",
         size=16, bold=True, color=WHITE)
add_footer(s, 7)

# ---------- SLIDE 8: MORE VERIFYWISE (governance & assurance) ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "The wider VerifyWise platform", kicker="More capabilities (1/2)")
add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5),
         "The LLM Gateway is one part. VerifyWise covers the full governance lifecycle.",
         size=17, color=DARK)

more1 = [
    ("AI systems & use cases", "Capture every AI request, run structured intake and score risk."),
    ("Model inventory", "Track all models across the organisation, with approval status."),
    ("Risk management & FRIA", "Risk register and fundamental rights impact assessments."),
    ("Compliance frameworks", "EU AI Act, ISO 42001, ISO 27001, NIST AI RMF, plus plugins."),
    ("Evidence hub", "Upload once, satisfy many controls. Audit-ready by default."),
    ("Approval workflows", "Multi-step sign-off chains with a full decision history."),
]
card_grid(s, more1, Inches(0.6), Inches(2.2), 3, Inches(4.0), Inches(1.55), Inches(0.15), Inches(0.2))
add_footer(s, 8)

# ---------- SLIDE 9: MORE VERIFYWISE (operate & disclose) ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "From oversight to public disclosure", kicker="More capabilities (2/2)")

more2 = [
    ("Vendor management", "Track third-party AI vendors and the models they supply."),
    ("Incident management", "Classify, route and report AI incidents, including Article 62."),
    ("Policy manager", "Author, version and assign AI policies across the organisation."),
    ("Shadow AI detection", "Find unsanctioned AI use before it becomes an incident."),
    ("AI trust center", "A public portal that shows customers and partners how you govern AI."),
    ("Tasks & automations", "Assign work and trigger actions when conditions are met."),
    ("Training registry", "Track AI literacy and training across teams."),
    ("Reporting & audit trail", "Every action logged, with reports auditors can use."),
]
card_grid(s, more2, Inches(0.6), Inches(1.5), 4, Inches(3.0), Inches(1.5), Inches(0.13), Inches(0.2),
          title_size=14, body_size=12)

add_rect(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.6), LIGHT, line=GREEN)
add_text(s, Inches(0.85), Inches(5.25), Inches(11.7), Inches(0.4),
         "ONE PLATFORM", size=12, bold=True, color=GREEN)
add_text(s, Inches(0.85), Inches(5.6), Inches(11.7), Inches(1.0),
         "16+ modules covering intake, risk, controls, evidence, audit and disclosure, across 24+ "
         "frameworks. Self-hosted or SaaS, available in Spanish, with our team alongside if you want it.",
         size=16, color=DARK)
add_footer(s, 9)

out = "/Users/gorkemcetin/Desktop/VerifyWise_AECC_EN.pptx"
prs.save(out)
print(out)
