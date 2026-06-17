from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

GREEN = RGBColor(0x13, 0x71, 0x5B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD5, 0xDD)
AMBER = RGBColor(0xB4, 0x6A, 0x00)

src = "/Users/gorkemcetin/Desktop/Casbah.pptx"
out = "/Users/gorkemcetin/Desktop/Casbah.pptx"

prs = Presentation(src)
SW, SH = prs.slide_width, prs.slide_height
BLANK = None
for layout in prs.slide_layouts:
    if "blank" in layout.name.lower() or layout.name == "BLANK":
        BLANK = layout
        break
if BLANK is None:
    BLANK = prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else -1]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "• " + it
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_header(slide, kicker, title):
    add_rect(slide, 0, 0, SW, Inches(0.08), GREEN)
    add_text(slide, Inches(0.6), Inches(0.25), Inches(10), Inches(0.4),
             kicker.upper(), size=12, bold=True, color=GREEN)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=DARK)
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(1.0), Inches(0.05), GREEN)


# ---------- SLIDE A: WHAT IS AN AI RISK? ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Workshop section 1", "What is an AI risk, and why is it different?")

add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.6),
         "An AI risk is a potential harm to people, the business or society arising from how an AI system is designed, trained, deployed or used.",
         size=15, bold=True, color=DARK)

# Four "why different" tiles
tiles = [
    ("Probabilistic, not deterministic",
        "Traditional software fails predictably. AI behaves differently on inputs it has never seen. Testing can't cover every case."),
    ("Emergent behaviour",
        "Risks appear after deployment that nobody designed in — bias, hallucinations, jailbreaks, drift. The system creates new risks during use."),
    ("Opacity",
        "You often can't explain why the model decided what it did. Traditional audit trails don't apply the same way."),
    ("Rapidly regulated",
        "EU AI Act, ISO 42001, NIST AI RMF, ICO guidance. The rules and expectations move faster than internal controls."),
]
x0, y0 = Inches(0.6), Inches(2.2)
tw, th = Inches(6.0), Inches(1.5)
gap_x, gap_y = Inches(0.1), Inches(0.15)
for i, (title, body) in enumerate(tiles):
    col = i % 2
    row = i // 2
    x = x0 + (tw + gap_x) * col
    y = y0 + (th + gap_y) * row
    add_rect(s, x, y, tw, th, LIGHT, line=BORDER)
    add_text(s, x + Inches(0.25), y + Inches(0.15), tw - Inches(0.35), Inches(0.4),
             title, size=15, bold=True, color=GREEN)
    add_text(s, x + Inches(0.25), y + Inches(0.55), tw - Inches(0.35), Inches(0.9),
             body, size=12, color=DARK)

# If not mitigated
add_rect(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.4), WHITE, line=GREEN)
add_text(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.4),
         "WHAT HAPPENS IF IT'S NOT MITIGATED", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.9), [
    "Regulatory: fines up to €35M or 7% of global turnover under the EU AI Act. ICO enforcement under UK GDPR.",
    "Reputational: one bad output going viral can erase years of brand work.   Legal: discrimination claims, breach of contract, data protection breaches.",
    "Operational: model has to be pulled, the business process behind it stops.   Trust: customers and regulators stop trusting the system entirely.",
], size=11)

# Footer
add_text(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
         "VerifyWise and The Casbah  |  Workshop", size=10, color=GREY)
add_text(s, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
         "Workshop", size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ---------- SLIDE B: EXAMPLES OF AI RISK + UK INCIDENT ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Workshop section 2", "Examples of AI risks — and one British case")

# Three reference frameworks side by side
add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.4),
         "The three reference catalogues practitioners use",
         size=14, bold=True, color=DARK)

refs = [
    ("MIT AI Risk Repository", "700+ documented risks across 7 domains. Academic, comprehensive. Use to discover."),
    ("IBM AI Risk Atlas", "About 70 plain-English risk categories. Practitioner-friendly. Use to categorise."),
    ("OWASP Top 10 for LLM Apps", "10 ranked, exploit-focused risks. Security-team friendly. Use to test."),
]
x = Inches(0.6)
w = Inches(4.0)
for title, body in refs:
    add_rect(s, x, Inches(1.85), w, Inches(1.4), LIGHT, line=BORDER)
    add_text(s, x + Inches(0.2), Inches(1.95), w - Inches(0.3), Inches(0.45),
             title, size=14, bold=True, color=GREEN)
    add_text(s, x + Inches(0.2), Inches(2.4), w - Inches(0.3), Inches(0.85),
             body, size=11, color=DARK)
    x += w + Inches(0.05)

# 5-category taxonomy
add_text(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.4),
         "A simple working taxonomy",
         size=14, bold=True, color=DARK)

cats = [
    ("Safety and accuracy", "Hallucination, drift, factually wrong output, unsafe action."),
    ("Bias and fairness", "Disparate treatment of protected groups, training-data bias."),
    ("Privacy and data", "PII leakage, memorisation of secrets, training on customer data."),
    ("Security", "Prompt injection, jailbreak, model theft, supply chain."),
    ("Governance", "Shadow AI, no human oversight, no audit trail, unclear ownership."),
]
x = Inches(0.6)
w = Inches(2.4)
gap = Inches(0.04)
for title, body in cats:
    add_rect(s, x, Inches(3.8), w, Inches(1.4), WHITE, line=BORDER)
    add_text(s, x + Inches(0.15), Inches(3.9), w - Inches(0.2), Inches(0.5),
             title, size=12, bold=True, color=GREEN)
    add_text(s, x + Inches(0.15), Inches(4.35), w - Inches(0.2), Inches(0.85),
             body, size=10, color=DARK)
    x += w + gap

# UK incident callout — Ofqual A-level scandal
add_rect(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.55), GREEN)
add_text(s, Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.4),
         "UK CASE STUDY — 2020 A-LEVEL ALGORITHM SCANDAL (OFQUAL)",
         size=11, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(5.78), Inches(11.7), Inches(0.55),
         "An algorithm regraded A-level results during COVID. 39% of grades were downgraded. Students from smaller, fee-paying schools benefited; those from larger state schools were systematically marked down.",
         size=12, color=WHITE)
add_text(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.5),
         "Maps to: Bias and fairness  |  Safety and accuracy  |  Governance.    Result: policy reversed within 4 days. Ofqual chair resigned.",
         size=12, bold=True, color=WHITE)

# Workshop hook
add_text(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
         "Ask the room: which of the five categories scares you most for your business?",
         size=10, bold=True, color=GREY)
add_text(s, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
         "Workshop", size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ---------- SLIDE C: WHO OWNS AI RISK? ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Workshop section 3", "Who owns AI risk in this organisation?")

add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
         "AI risk has no natural home. That is the problem — and the conversation we need to have today.",
         size=15, bold=True, color=DARK)

# RACI-style grid: 4 functions, what they think they own
owners = [
    ("IT and Security",
        "Owns the infrastructure and the access controls.",
        "But doesn't decide what the AI does, or who it affects."),
    ("Legal and Compliance",
        "Owns the regulatory interpretation.",
        "But can't fix the model, the training data or the deployment."),
    ("The business unit using the AI",
        "Owns the outcome and the customer relationship.",
        "But often didn't choose the model and can't audit it."),
    ("Executive leadership",
        "Owns accountability if something goes wrong.",
        "But is usually furthest from the actual system."),
]
x0, y0 = Inches(0.6), Inches(2.15)
tw, th = Inches(6.0), Inches(1.4)
gap_x, gap_y = Inches(0.1), Inches(0.15)
for i, (who, owns, but) in enumerate(owners):
    col = i % 2
    row = i // 2
    x = x0 + (tw + gap_x) * col
    y = y0 + (th + gap_y) * row
    add_rect(s, x, y, tw, th, LIGHT, line=BORDER)
    add_text(s, x + Inches(0.25), y + Inches(0.12), tw - Inches(0.35), Inches(0.4),
             who, size=15, bold=True, color=GREEN)
    add_text(s, x + Inches(0.25), y + Inches(0.5), tw - Inches(0.35), Inches(0.4),
             "Owns: " + owns, size=11, bold=True, color=DARK)
    add_text(s, x + Inches(0.25), y + Inches(0.88), tw - Inches(0.35), Inches(0.45),
             "Gap: " + but, size=11, color=GREY)

# Bottom takeaway bar
add_rect(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.5), WHITE, line=GREEN)
add_text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.4),
         "THE QUESTION FOR THE ROOM", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(1.0), [
    "Who in your organisation gets the phone call when an AI system causes harm?",
    "Is that person empowered to do anything about it?",
    "If your answer is more than one name, you don't have an owner — you have a committee.",
], size=12)

# Footer
add_text(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
         "Goal of this section: leave with one named owner for AI risk in this organisation.",
         size=10, bold=True, color=GREY)
add_text(s, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
         "Workshop", size=10, color=GREY, align=PP_ALIGN.RIGHT)


prs.save(out)
print(out)
