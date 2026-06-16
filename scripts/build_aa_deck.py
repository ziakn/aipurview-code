from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

GREEN = RGBColor(0x13, 0x71, 0x5B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD5, 0xDD)
AMBER = RGBColor(0xB4, 0x6A, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

TOTAL = 16


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "• " + it
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_header(slide, title, kicker=None):
    add_rect(slide, 0, 0, SW, Inches(0.08), GREEN)
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.22), Inches(10), Inches(0.4),
                 kicker.upper(), size=11, bold=True, color=GREEN)
        add_text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
                 title, size=26, bold=True, color=DARK)
    else:
        add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                 title, size=26, bold=True, color=DARK)
    add_rect(slide, Inches(0.6), Inches(1.15), Inches(1.0), Inches(0.05), GREEN)


def add_footer(slide, page):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(10), Inches(0.3),
             "VerifyWise Ltd  |  SEIS/EIS Advanced Assurance Application  |  Confidential", size=9, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{page} / {TOTAL}", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def placeholder(slide, x, y, w, h, label):
    add_rect(slide, x, y, w, h, LIGHT, line=AMBER)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.3),
             "TO COMPLETE", size=9, bold=True, color=AMBER)
    add_text(slide, x + Inches(0.2), y + Inches(0.35), w - Inches(0.4), h - Inches(0.4),
             label, size=12, color=DARK)


# ---------- SLIDE 1: COVER ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, Inches(4.6), SH, GREEN)
add_text(s, Inches(0.6), Inches(0.6), Inches(3.5), Inches(0.5),
         "VERIFYWISE", size=14, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.8), Inches(3.8), Inches(1.4),
         "Application for\nAdvanced\nAssurance.", size=32, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(5.0), Inches(3.8), Inches(0.5),
         "SEIS / EIS", size=16, color=WHITE)

add_text(s, Inches(5.2), Inches(0.8), Inches(7.5), Inches(0.4),
         "HMRC SEIS/EIS ADVANCED ASSURANCE", size=12, bold=True, color=GREEN)
add_text(s, Inches(5.2), Inches(1.2), Inches(7.5), Inches(2.2),
         "VerifyWise Ltd —\nbusiness plan and\nstatutory declarations.",
         size=32, bold=True, color=DARK)
add_rect(s, Inches(5.2), Inches(3.9), Inches(0.8), Inches(0.05), GREEN)
add_text(s, Inches(5.2), Inches(4.05), Inches(7.5), Inches(0.4),
         "Submitted in support of HMRC Advanced Assurance application", size=14, bold=True, color=DARK)
add_text(s, Inches(5.2), Inches(4.45), Inches(7.5), Inches(0.4),
         "Applicant: VerifyWise Ltd  |  Company number 17045757", size=13, color=GREY)
add_text(s, Inches(5.2), Inches(4.75), Inches(7.5), Inches(0.4),
         "Signed by: Dr Gorkem Cetin, Director", size=13, color=GREY)
add_text(s, Inches(5.2), Inches(5.05), Inches(7.5), Inches(0.4),
         "Date: June 2026", size=13, color=GREY)

add_rect(s, Inches(5.2), Inches(5.8), Inches(7.5), Inches(1.1), LIGHT, line=BORDER)
add_text(s, Inches(5.4), Inches(5.9), Inches(7.2), Inches(0.35),
         "REGISTERED OFFICE", size=10, bold=True, color=GREEN)
add_text(s, Inches(5.4), Inches(6.2), Inches(7.2), Inches(0.65),
         "Suite 101, Lumina Business Centre,\n32 Lumina Way, London, England, EN1 1FS",
         size=12, color=DARK)

# ---------- SLIDE 2: COMPANY FACTS ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Company facts", kicker="Section 1")

rows = [
    ("Registered name", "VERIFYWISE LTD"),
    ("Company registration number", "17045757"),
    ("Date of incorporation", "20 February 2026"),
    ("Registered office", "Suite 101, Lumina Business Centre, 32 Lumina Way, London, England, EN1 1FS"),
    ("Country of incorporation", "England and Wales"),
    ("Company type", "Private limited company"),
    ("Company status", "Active"),
    ("SIC codes", "58290 (Other software publishing); 62012 (Business and domestic software development)"),
    ("Director", "Dr Gorkem Cetin — appointed 20 February 2026 (sole director on the register)"),
    ("Person with significant control", "Dr Gorkem Cetin (per Companies House). PSC filing to be updated to reflect current shareholding (Dr Gorkem Cetin and Mr Ulas Ozguven)."),
    ("Other shareholders / officers", "Mr Ulas Ozguven — shareholder and manager (not a statutory director)."),
    ("Accounts reference date", "First accounts to 28 February 2027 (due 20 November 2027)"),
    ("Confirmation statement", "First statement date 19 February 2027 (due 5 March 2027)"),
]

y = Inches(1.4)
for label, val in rows:
    add_rect(s, Inches(0.6), y, Inches(4.2), Inches(0.42), LIGHT, line=BORDER)
    add_text(s, Inches(0.75), y + Inches(0.08), Inches(4.0), Inches(0.3),
             label, size=11, bold=True, color=DARK)
    add_rect(s, Inches(4.8), y, Inches(7.9), Inches(0.42), WHITE, line=BORDER)
    add_text(s, Inches(4.95), y + Inches(0.08), Inches(7.7), Inches(0.3),
             val, size=11, color=DARK)
    y += Inches(0.44)

add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.25),
         "Source: Companies House public register, retrieved June 2026.",
         size=9, color=GREY)
add_footer(s, 2)

# ---------- SLIDE 3: QUALIFYING TRADE ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Qualifying trade statement", kicker="Section 2 — Statutory")

add_text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5),
         "VerifyWise's trade is a qualifying trade for SEIS and EIS purposes.",
         size=18, bold=True, color=DARK)

add_rect(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.0), LIGHT, line=BORDER)
add_text(s, Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.4),
         "WHAT VERIFYWISE DOES", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.4),
         "VerifyWise Ltd develops and licenses an AI governance software platform. Revenue comes from SaaS "
         "subscriptions, enterprise licences for self-hosted deployments, and professional services. The company "
         "owns all platform intellectual property. R&D is currently delivered through founder time and contracted "
         "engineering capacity based in Canada; the funds raised in this round will build out UK engineering and "
         "commercial functions. The trade is a software development and licensing trade carried on by the UK company.",
         size=13, color=DARK)

add_rect(s, Inches(0.6), Inches(4.15), Inches(12.1), Inches(2.5), WHITE, line=BORDER)
add_text(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.4),
         "EXCLUDED ACTIVITIES — NONE APPLY", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(4.65), Inches(11.7), Inches(1.9), [
    "No dealing in land, commodities, futures, shares, securities or other financial instruments.",
    "No banking, insurance, money-lending, debt-factoring, hire-purchase financing or other financial activities.",
    "No leasing, letting ships on charter or receiving royalties or licence fees as a principal trade — IP is owned and exploited as part of the trading activity, not licensed out passively.",
    "No legal or accountancy services. No property development. No farming, forestry, hotels, nursing or residential care homes.",
    "No coal or steel production, no shipbuilding, no electricity generation other than incidental.",
])
add_text(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.3),
         "References: ITA 2007 s.189–199 (qualifying trades and excluded activities).",
         size=10, color=GREY)
add_footer(s, 3)

# ---------- SLIDE 4: PROBLEM AND MARKET ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "The problem and the market", kicker="Section 3")

add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
         "AI regulation is creating a new compliance category. The market is opening now.",
         size=18, bold=True, color=DARK)

tiles = [
    ("EU AI Act", "Obligations phased through 2026 and 2027. Maximum fine is the higher of €35M or 7% of global turnover."),
    ("ISO 42001", "First international AI management-system standard. Published December 2023. Certification audits are live."),
    ("CEN-CENELEC", "JTC 21 harmonised standards expected Q4 2026. prEN 18286 already in public enquiry — quality management for the AI Act."),
]
x = Inches(0.6)
for title, body in tiles:
    add_rect(s, x, Inches(2.1), Inches(4.0), Inches(2.4), LIGHT, line=BORDER)
    add_text(s, x + Inches(0.25), Inches(2.25), Inches(3.6), Inches(0.5),
             title, size=20, bold=True, color=GREEN)
    add_text(s, x + Inches(0.25), Inches(2.85), Inches(3.6), Inches(1.5),
             body, size=13, color=DARK)
    x += Inches(4.2)

add_rect(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.9), WHITE, line=BORDER)
add_text(s, Inches(0.8), Inches(4.95), Inches(11.7), Inches(0.4),
         "MARKET SIZE", size=11, bold=True, color=GREEN)
placeholder(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.3),
            "TAM/SAM/SOM figures with Gartner, IDC or Forrester citations. Include AI governance, GRC and compliance software market sizes with 3-5 year CAGR.")
add_footer(s, 4)

# ---------- SLIDE 5: PRODUCT ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "The product", kicker="Section 4")

add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
         "VerifyWise is the open-source AI governance platform for the regulatory era.",
         size=18, bold=True, color=DARK)

caps = [
    ("AI system inventory", "Every model, vendor and use case in one register."),
    ("Multi-framework engine", "EU AI Act, ISO 42001, ISO 27001, NIST AI RMF plus 30+ plugins."),
    ("Risk register and FRIA", "Continuous risk management aligned to EU AI Act Article 9."),
    ("Evidence hub", "Upload once, satisfy many controls. Audit-ready."),
    ("AI Gateway and guardrails", "PII detection, spend caps, policy checks at runtime."),
    ("Shadow AI detection", "Find unsanctioned AI use across the enterprise."),
]
x0, y0 = Inches(0.6), Inches(2.0)
tw, th = Inches(4.0), Inches(1.4)
gap_x, gap_y = Inches(0.15), Inches(0.2)
for i, (title, body) in enumerate(caps):
    col = i % 3
    row = i // 3
    x = x0 + (tw + gap_x) * col
    y = y0 + (th + gap_y) * row
    add_rect(s, x, y, tw, th, WHITE, line=BORDER)
    add_text(s, x + Inches(0.25), y + Inches(0.15), tw - Inches(0.35), Inches(0.4),
             title, size=14, bold=True, color=DARK)
    add_text(s, x + Inches(0.25), y + Inches(0.55), tw - Inches(0.35), Inches(0.8),
             body, size=12, color=GREY)

add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
         "Open-core model. Self-hosted or SaaS. Multi-tenant. Built for regulated industries.",
         size=14, bold=True, color=GREEN)
placeholder(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.8),
            "Insert product screenshot or architecture diagram before submission.")
add_footer(s, 5)

# ---------- SLIDE 6: BUSINESS MODEL ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Business model and revenue", kicker="Section 5")

streams = [
    ("SaaS subscriptions", "Per-seat or per-tenant. Recurring monthly or annual."),
    ("Enterprise licence", "Self-hosted deployments. Annual contract."),
    ("Professional services", "Implementation, integration, framework configuration."),
    ("Training and certification", "Customer enablement programmes. Margin-positive."),
]
x = Inches(0.6)
w = Inches(2.95)
for title, body in streams:
    add_rect(s, x, Inches(1.4), w, Inches(2.4), LIGHT, line=BORDER)
    add_text(s, x + Inches(0.2), Inches(1.55), w - Inches(0.3), Inches(0.5),
             title, size=15, bold=True, color=GREEN)
    add_text(s, x + Inches(0.2), Inches(2.1), w - Inches(0.3), Inches(1.6),
             body, size=12, color=DARK)
    x += w + Inches(0.15)

add_rect(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(2.7), WHITE, line=BORDER)
add_text(s, Inches(0.8), Inches(4.15), Inches(11.7), Inches(0.4),
         "CURRENT PRICING AND UNIT ECONOMICS", size=11, bold=True, color=GREEN)
placeholder(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.1),
            "Insert: average contract value, gross margin %, CAC, LTV, payback period, current ARR and projected ARR at end of year 1/2/3.")
add_footer(s, 6)

# ---------- SLIDE 7: TRACTION ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Traction to date", kicker="Section 6")

tiles = [
    ("Revenue", "[Current MRR / ARR]\nLast 12 months: [£]\nGrowth: [%]"),
    ("Customers and pilots", "[# paying customers]\n[# active pilots]\nNamed logos: [list]"),
    ("Open-source adoption", "GitHub stars: [#]\nContributors: [#]\nDeployments: [estimated]"),
    ("Partnerships", "Responsible AI Institute\n[Other named partners]\n[SI relationships]"),
]
x0, y0 = Inches(0.6), Inches(1.4)
tw, th = Inches(3.0), Inches(2.6)
gap = Inches(0.1)
for i, (title, body) in enumerate(tiles):
    x = x0 + (tw + gap) * i
    add_rect(s, x, y0, tw, th, LIGHT, line=BORDER)
    add_text(s, x + Inches(0.2), y0 + Inches(0.15), tw - Inches(0.3), Inches(0.5),
             title, size=15, bold=True, color=GREEN)
    add_text(s, x + Inches(0.2), y0 + Inches(0.7), tw - Inches(0.3), Inches(1.8),
             body, size=13, color=DARK)

add_rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5), WHITE, line=BORDER)
add_text(s, Inches(0.8), Inches(4.35), Inches(11.7), Inches(0.4),
         "PROOF POINTS", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.9), [
    "[Revenue milestone — e.g. first £X of ARR achieved in Q# 202#]",
    "[Notable customer or design partner won, with sector and use case]",
    "[Open-source community milestone — stars threshold, key external contributors]",
    "[Inbound demand signal — pipeline value, demo requests/month]",
    "[Industry recognition — Gartner mentions, analyst briefings, awards]",
])
add_footer(s, 7)

# ---------- SLIDE 8: TEAM ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Team and capability", kicker="Section 7")

add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
         "Founder-led, with a roadmap to build out a UK-based engineering and commercial team.",
         size=15, bold=True, color=DARK)

members = [
    ("Dr Gorkem Cetin", "Director, founding shareholder",
        "Appointed 20 February 2026 on incorporation. Companies House identity-verified. "
        "Country of residence: Canada. Leads product, R&D and overall company direction."),
    ("Mr Ulas Ozguven", "Manager and shareholder (not a statutory director)",
        "Holds a minority shareholding in VerifyWise Ltd. Manages day-to-day commercial activities. "
        "PSC and shareholder register to be updated at Companies House to reflect his shareholding."),
]
y = Inches(1.95)
for name, role, bio in members:
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.15), WHITE, line=BORDER)
    add_text(s, Inches(0.9), y + Inches(0.12), Inches(5.0), Inches(0.4),
             name, size=16, bold=True, color=DARK)
    add_text(s, Inches(0.9), y + Inches(0.55), Inches(5.0), Inches(0.5),
             role, size=11, bold=True, color=GREEN)
    add_text(s, Inches(6.0), y + Inches(0.15), Inches(6.0), Inches(0.95),
             bio, size=11, color=DARK)
    y += Inches(1.25)

add_rect(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.35), LIGHT, line=BORDER)
add_text(s, Inches(0.8), Inches(4.65), Inches(11.7), Inches(0.4),
         "PLANNED TEAM BUILD-OUT (FUNDED BY THIS ROUND)", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.85), [
    "UK-based engineering hires to deliver EU AI Act, ISO 42001 and CEN-CENELEC modules to GA.",
    "UK-based commercial lead for enterprise sales and SI partnerships.",
    "UK-based customer success lead for paying enterprise accounts.",
    "Advisors: [to be appointed — AI policy advisor, EIS-experienced non-executive director, regulated-industry buyer advisor].",
], size=12)

add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.25),
         "Director and shareholder details verified against Companies House public register, June 2026.",
         size=9, color=GREY)
add_footer(s, 8)

# ---------- SLIDE 9: GROUP STRUCTURE, IP AND UK PE ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Group structure, IP ownership and UK trade", kicker="Section 8 — Statutory")

add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "VerifyWise Ltd is an independent UK company. Trade is carried on in the UK. All IP is owned by VerifyWise Ltd.",
         size=15, bold=True, color=DARK)

add_rect(s, Inches(0.6), Inches(1.9), Inches(6.0), Inches(2.4), LIGHT, line=BORDER)
add_text(s, Inches(0.8), Inches(2.05), Inches(5.7), Inches(0.4),
         "INDEPENDENCE", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(2.4), Inches(5.7), Inches(1.9), [
    "VerifyWise Ltd is not a subsidiary of any other company.",
    "No parent company exists in the UK or overseas.",
    "Sole PSC: Dr Gorkem Cetin (75%+ shares and voting rights), per Companies House.",
    "No 50%+ ownership by another company. Independence test met.",
], size=12)

add_rect(s, Inches(6.75), Inches(1.9), Inches(5.95), Inches(2.4), WHITE, line=BORDER)
add_text(s, Inches(6.95), Inches(2.05), Inches(5.65), Inches(0.4),
         "INTELLECTUAL PROPERTY", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(6.95), Inches(2.4), Inches(5.65), Inches(1.9), [
    "All platform IP — source code, trademarks, documentation — is owned by VerifyWise Ltd.",
    "No licensing of IP from a foreign parent or related party.",
    "Open-source components used under standard permissive licences; do not affect VerifyWise's ownership of its own work.",
    "IP licensing is not the company's principal trade — software development and sales is.",
], size=12)

add_rect(s, Inches(0.6), Inches(4.45), Inches(12.1), Inches(2.4), LIGHT, line=BORDER)
add_text(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.4),
         "UK PERMANENT ESTABLISHMENT", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(4.95), Inches(11.7), Inches(1.9), [
    "Registered office at Suite 101, Lumina Business Centre, 32 Lumina Way, London, EN1 1FS.",
    "Contracts with customers, suppliers and contractors are entered into by VerifyWise Ltd, a UK company. Revenue, expenditure and statutory accounts are recorded in the UK.",
    "Sole director Dr Gorkem Cetin is currently resident in Canada. R&D is presently delivered by contracted engineering capacity based in Canada under contract to VerifyWise Ltd.",
    "Proceeds from this round will fund UK-based hires (engineering and commercial). UK PAYE will be operated from the point the first UK employee starts.",
    "VerifyWise Ltd is registered for UK corporation tax. Board meetings, decision-making and shareholder records are kept by the UK company.",
], size=12)

add_footer(s, 9)

# ---------- SLIDE 10: KIC ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Knowledge-intensive company case", kicker="Section 9 — Statutory")

add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
         "VerifyWise meets the Knowledge-Intensive Company (KIC) conditions under ITA 2007 s.252A.",
         size=16, bold=True, color=DARK)

add_rect(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.6), LIGHT, line=BORDER)
add_text(s, Inches(0.8), Inches(2.15), Inches(5.7), Inches(0.4),
         "OPERATING COSTS CONDITION", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.8), Inches(2.5), Inches(5.7), Inches(0.4),
         "At least one of the following:", size=13, bold=True, color=DARK)
add_bullets(s, Inches(0.8), Inches(2.95), Inches(5.7), Inches(2.0), [
    "≥15% of operating costs spent on R&D in one of the previous 3 years, OR",
    "≥10% of operating costs spent on R&D in each of the previous 3 years.",
], size=12)
placeholder(s, Inches(0.8), Inches(4.95), Inches(5.7), Inches(1.55),
            "Insert R&D spend as % of operating costs for each of the last 3 financial years.")

add_rect(s, Inches(6.75), Inches(2.0), Inches(5.95), Inches(4.6), WHITE, line=BORDER)
add_text(s, Inches(6.95), Inches(2.15), Inches(5.65), Inches(0.4),
         "INNOVATION OR SKILLED EMPLOYEE CONDITION", size=11, bold=True, color=GREEN)
add_text(s, Inches(6.95), Inches(2.5), Inches(5.65), Inches(0.4),
         "VerifyWise meets the innovation condition:", size=13, bold=True, color=DARK)
add_bullets(s, Inches(6.95), Inches(2.95), Inches(5.65), Inches(3.5), [
    "Company creates intellectual property in the form of proprietary AI governance software.",
    "Within 10 years, the majority of trade is expected to consist of business arising from that IP.",
    "Activities are not routine — they involve ongoing applied research into AI risk classification, evidence reuse across regulatory frameworks, and runtime guardrails.",
    "Alternatively, the skilled-employee condition is met if ≥20% of full-time-equivalent employees hold master's-level or higher qualifications and are engaged in R&D.",
], size=12)

add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3),
         "If KIC status is confirmed: EIS lifetime cap rises to £20M, age limit to 10 years, employee cap to 500.",
         size=11, color=GREY)
add_footer(s, 10)

# ---------- SLIDE 11: GROWTH PLAN ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Growth plan — three years", kicker="Section 10")

years = [
    ("Year 1", "2026-27", [
        "Ship EU AI Act and ISO 42001 modules to GA",
        "[#] paying enterprise customers",
        "[Headcount] full-time staff in the UK",
        "ARR target: £[X]",
    ]),
    ("Year 2", "2027-28", [
        "NIST AI RMF and CEN-CENELEC harmonised standards coverage",
        "[#] enterprise customers, including [#] outside the UK",
        "Channel partnerships with [#] SIs (NTT DATA, [others])",
        "ARR target: £[X]",
    ]),
    ("Year 3", "2028-29", [
        "Auditor and consultant ecosystem live",
        "International expansion — EU and North America",
        "[Headcount] full-time staff",
        "ARR target: £[X]",
    ]),
]
x = Inches(0.6)
w = Inches(4.0)
for label, period, items in years:
    add_rect(s, x, Inches(1.4), w, Inches(5.2), WHITE, line=BORDER)
    add_rect(s, x, Inches(1.4), w, Inches(0.7), GREEN)
    add_text(s, x + Inches(0.2), Inches(1.5), w - Inches(0.3), Inches(0.4),
             label, size=18, bold=True, color=WHITE)
    add_text(s, x + Inches(0.2), Inches(1.85), w - Inches(0.3), Inches(0.3),
             period, size=11, bold=True, color=WHITE)
    add_bullets(s, x + Inches(0.2), Inches(2.3), w - Inches(0.3), Inches(4.1),
                items, size=13)
    x += w + Inches(0.15)

add_footer(s, 11)

# ---------- SLIDE 12: USE OF FUNDS ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Use of investment proceeds", kicker="Section 11")

add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
         "Round size: £200,000 for 10% of the company. Funds applied to growth of the qualifying trade.",
         size=16, bold=True, color=DARK)

rows = [
    ("Engineering", "50%", "£100,000 — UK-based engineering hires to deliver EU AI Act, ISO 42001 and CEN-CENELEC modules to GA."),
    ("Sales and partnerships", "20%", "£40,000 — UK commercial lead, enterprise account development and SI partnerships (NTT DATA and others)."),
    ("Marketing and demand generation", "15%", "£30,000 — content, events, paid acquisition and analyst relations."),
    ("Infrastructure and security", "10%", "£20,000 — cloud, observability, SOC 2 readiness, penetration testing."),
    ("Legal and compliance", "5%", "£10,000 — counsel for contracts, IP protection and regulatory advice."),
]
y = Inches(2.0)
for label, pct, detail in rows:
    add_rect(s, Inches(0.6), y, Inches(4.0), Inches(0.7), LIGHT, line=BORDER)
    add_text(s, Inches(0.8), y + Inches(0.2), Inches(2.8), Inches(0.4),
             label, size=14, bold=True, color=DARK)
    add_text(s, Inches(3.4), y + Inches(0.2), Inches(1.2), Inches(0.4),
             pct, size=16, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    add_rect(s, Inches(4.6), y, Inches(8.1), Inches(0.7), WHITE, line=BORDER)
    add_text(s, Inches(4.8), y + Inches(0.2), Inches(7.7), Inches(0.4),
             detail, size=12, color=DARK)
    y += Inches(0.75)

add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
         "All funds will be spent within 24 months of issue (EIS) / 36 months (SEIS) and applied wholly to growth of the qualifying trade.",
         size=11, bold=True, color=GREEN)
add_footer(s, 12)

# ---------- SLIDE 13: RISK TO CAPITAL ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Risk-to-capital condition", kicker="Section 12 — Statutory")

add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "The investment meets the risk-to-capital condition under ITA 2007 s.157A.",
         size=16, bold=True, color=DARK)

add_rect(s, Inches(0.6), Inches(1.9), Inches(6.0), Inches(4.8), LIGHT, line=BORDER)
add_text(s, Inches(0.8), Inches(2.05), Inches(5.7), Inches(0.4),
         "LONG-TERM GROWTH OBJECTIVE", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(2.45), Inches(5.7), Inches(4.0), [
    "VerifyWise has a clear plan to grow revenue, headcount and customer base over the next 3-5 years.",
    "Proceeds fund engineering, sales and marketing — not shareholder distributions.",
    "Growth is geographic (UK to EU to North America), product (modules, integrations) and channel (SI partnerships).",
    "All employees hired with this round will be in the UK or will support UK trade.",
], size=12)

add_rect(s, Inches(6.75), Inches(1.9), Inches(5.95), Inches(4.8), WHITE, line=BORDER)
add_text(s, Inches(6.95), Inches(2.05), Inches(5.65), Inches(0.4),
         "GENUINE RISK OF LOSS OF CAPITAL", size=11, bold=True, color=GREEN)
add_bullets(s, Inches(6.95), Inches(2.45), Inches(5.65), Inches(4.0), [
    "Early-stage software company. Pre-profit. No guaranteed return.",
    "No preferential rights on the shares being issued — ordinary shares only.",
    "No arrangements for capital protection, guaranteed exit or share buy-back.",
    "Investors face significant risk that the company may not succeed and that capital may be lost in full.",
    "Market risk: AI governance is an emerging category and customer adoption timelines are not certain.",
], size=12)

add_footer(s, 13)

# ---------- SLIDE 14: SHARE STRUCTURE ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Share structure and round terms", kicker="Section 13 — Statutory")

rows = [
    ("Type of shares to be issued", "Ordinary shares — no preferential rights to dividends, assets on winding up or redemption."),
    ("Shares paid up in cash", "Yes. Shares fully subscribed and paid up in cash on issue."),
    ("Round size", "£200,000"),
    ("Equity offered", "10% of the company (post-money)"),
    ("Pre-money valuation", "£1,800,000"),
    ("Post-money valuation", "£2,000,000"),
    ("Share price per share", "To be set at allotment so that the round delivers 10% on a fully diluted basis."),
    ("Dilution to existing shareholders", "10%"),
    ("Anchor investor", "WMtech Partners (International) Ltd — 71-75 Shelton Street, London, WC2H 9JQ (corporate subscriber; SEIS/EIS relief does not apply to WMtech's subscription)."),
    ("Investor contact", "Moqsood Ali, Partner — moqsood@wmtech.io  |  +44 (0) 7817 420 061  |  wmtech.io"),
    ("Individual SEIS/EIS subscribers", "Round structured to admit individual SEIS/EIS-eligible investors alongside WMtech. Named individual subscribers will be confirmed before share issue."),
    ("Capital protection arrangements", "None. No put options, no guaranteed buy-back, no preference."),
    ("Investor exit arrangements", "None pre-agreed. Liquidity expected via trade sale or future financing rounds at market valuation."),
    ("Use of money condition", "Funds applied to qualifying trade within 24 months of share issue (EIS) or 3 years (SEIS)."),
]
y = Inches(1.35)
for label, val in rows:
    add_rect(s, Inches(0.6), y, Inches(4.6), Inches(0.38), LIGHT, line=BORDER)
    add_text(s, Inches(0.75), y + Inches(0.06), Inches(4.4), Inches(0.3),
             label, size=10, bold=True, color=DARK)
    add_rect(s, Inches(5.2), y, Inches(7.5), Inches(0.38), WHITE, line=BORDER)
    add_text(s, Inches(5.35), y + Inches(0.06), Inches(7.3), Inches(0.3),
             val, size=10, color=DARK)
    y += Inches(0.4)

add_footer(s, 14)

# ---------- SLIDE 15: STATUTORY CHECKLIST ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Statutory compliance checklist", kicker="Section 14 — Statutory")

left = [
    ("Incorporation date", "20 February 2026 — well within SEIS (<3y), EIS (<7y) and KIC (<10y) limits"),
    ("Gross assets pre-investment", "£[X] — to confirm; well below SEIS £350k / EIS £15M limits"),
    ("Gross assets post-investment", "£[X] — to confirm; below SEIS £350k / EIS £16M limit"),
    ("Full-time-equivalent employees", "[#] — well below SEIS 25 / EIS 250 / KIC 500 limit"),
    ("UK permanent establishment", "Yes — see Section 8"),
    ("Independent (not >50% owned by another company)", "Yes — sole PSC is an individual (Companies House)"),
]
right = [
    ("Trading status at investment date", "Trading or preparing to trade"),
    ("Funds raised under SEIS to date", "£[X] (SEIS lifetime cap £250k)"),
    ("Funds raised under EIS / VCT to date", "£[X] (EIS lifetime cap £12M / KIC £20M)"),
    ("Annual investment cap (last 12 months)", "£[X] (£5M standard / £10M KIC)"),
    ("Shares will be held for 3 years minimum", "Confirmed in shareholders' agreement"),
    ("No linked loans, guarantees or repayments to investors", "Confirmed"),
]

add_text(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.4),
         "Eligibility", size=14, bold=True, color=GREEN)
y = Inches(1.85)
for label, val in left:
    add_rect(s, Inches(0.6), y, Inches(6.0), Inches(0.75), LIGHT, line=BORDER)
    add_text(s, Inches(0.75), y + Inches(0.08), Inches(5.7), Inches(0.3),
             label, size=11, bold=True, color=DARK)
    add_text(s, Inches(0.75), y + Inches(0.38), Inches(5.7), Inches(0.35),
             val, size=11, color=GREY)
    y += Inches(0.78)

add_text(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(0.4),
         "Investment integrity", size=14, bold=True, color=GREEN)
y = Inches(1.85)
for label, val in right:
    add_rect(s, Inches(6.85), y, Inches(5.85), Inches(0.75), WHITE, line=BORDER)
    add_text(s, Inches(7.0), y + Inches(0.08), Inches(5.55), Inches(0.3),
             label, size=11, bold=True, color=DARK)
    add_text(s, Inches(7.0), y + Inches(0.38), Inches(5.55), Inches(0.35),
             val, size=11, color=GREY)
    y += Inches(0.78)

add_footer(s, 15)

# ---------- SLIDE 16: APPENDIX ----------
s = prs.slides.add_slide(BLANK)
add_header(s, "Appendix and supporting documents", kicker="Section 15")

docs = [
    ("Certificate of incorporation", "Companies House — attached separately."),
    ("Memorandum and articles of association", "Latest filed version — attached separately."),
    ("Latest filed accounts", "[Year] statutory accounts — attached separately."),
    ("Shareholders' agreement / term sheet", "Draft for proposed round — attached separately."),
    ("Three-year financial forecast", "Profit and loss, cash flow, headcount — spreadsheet attached separately."),
    ("R&D spend evidence", "Breakdown supporting the KIC operating-costs condition — spreadsheet attached separately."),
    ("Director and shareholder list", "Per Companies House filings — attached separately."),
    ("Letter from anchor investor", "WMtech Partners (International) Ltd — letter of intent signed by Moqsood Ali, Partner, attached separately."),
]
y = Inches(1.4)
for label, val in docs:
    add_rect(s, Inches(0.6), y, Inches(4.6), Inches(0.6), LIGHT, line=BORDER)
    add_text(s, Inches(0.75), y + Inches(0.15), Inches(4.4), Inches(0.3),
             label, size=13, bold=True, color=DARK)
    add_rect(s, Inches(5.2), y, Inches(7.5), Inches(0.6), WHITE, line=BORDER)
    add_text(s, Inches(5.35), y + Inches(0.15), Inches(7.3), Inches(0.3),
             val, size=12, color=DARK)
    y += Inches(0.63)

add_rect(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5), GREEN)
add_text(s, Inches(0.85), Inches(6.6), Inches(11.7), Inches(0.4),
         "Contact for queries: Dr Gorkem Cetin, Director, VerifyWise Ltd — gorkem@verifywise.ai",
         size=13, bold=True, color=WHITE)
add_footer(s, 16)

out = "/Users/gorkemcetin/Desktop/VerifyWise_Advanced_Assurance.pptx"
prs.save(out)
print(out)
